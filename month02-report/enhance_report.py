from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).parent
INPUT = ROOT / "month02-text.pptx"
OUTPUT = ROOT / "Topic02_Month02_Report.pptx"
WEB_SCREENSHOT = ROOT / "assets" / "web-ui-ja.png"

NAVY = RGBColor(0x00, 0x37, 0x6F)
RED = RGBColor(0xC0, 0x00, 0x00)
CYAN = RGBColor(0xD9, 0xF1, 0xF4)
PALE_BLUE = RGBColor(0xE9, 0xF3, 0xFA)
PALE_GREEN = RGBColor(0xE7, 0xF4, 0xE4)
GRAY = RGBColor(0x5B, 0x67, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_cell_text(cell, text, bold=False):
    cell.text = text
    frame = cell.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in frame.paragraphs:
        paragraph.font.name = "Times New Roman"
        paragraph.font.size = Pt(17)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Times New Roman"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, step, title, body, fill, active=False):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = RED if active else NAVY
    card.line.width = Pt(2.5 if active else 1.2)

    add_text(slide, x + 0.18, y + 0.22, w - 0.36, 0.38, step, 14, RED if active else NAVY, True)
    add_text(slide, x + 0.18, y + 0.72, w - 0.36, 0.85, title, 19, NAVY, True)
    add_text(slide, x + 0.2, y + 1.7, w - 0.4, 1.55, body, 13, GRAY)


prs = Presentation(INPUT)

# Cover metadata
table = prs.slides[0].shapes[5].table
cover_rows = [
    ("Month", "Month 2 / From 01/06/2026 - 01/07/2026"),
    ("Member", "AKIRA, SHUICHI"),
    ("Research Theme", "Local LLM-Assisted Railway Terminology Extraction"),
    ("Mentor", "Dr. BUI QUOC BAO"),
]
for row, (label, value) in enumerate(cover_rows):
    set_cell_text(table.cell(row, 0), label, True)
    set_cell_text(table.cell(row, 1), value, row == 2)

# Replace the obsolete BERT table screenshot with the current multilingual web UI.
web_slide = prs.slides[8]
cover = web_slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.25), Inches(2.18), Inches(18.1), Inches(6.72)
)
cover.fill.solid()
cover.fill.fore_color.rgb = WHITE
cover.line.color.rgb = RGBColor(0xC8, 0xD4, 0xDC)
cover.line.width = Pt(1.2)

web_slide.shapes.add_picture(
    str(WEB_SCREENSHOT), Inches(4.37), Inches(2.34), width=Inches(11.26), height=Inches(6.41)
)
add_text(web_slide, 15.75, 2.5, 3.2, 0.72, "日本語 / EN / VI", 17, NAVY, True)

# Replace the outdated calendar graphic with the actual inference-to-training roadmap.
roadmap_slide = prs.slides[10]
roadmap_cover = roadmap_slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.18), Inches(2.72), Inches(17.72), Inches(6.7)
)
roadmap_cover.fill.solid()
roadmap_cover.fill.fore_color.rgb = WHITE
roadmap_cover.line.color.rgb = WHITE

steps = [
    ("01 · CURRENT", "Inference baseline", "Run Qwen locally\nand export draft CSV.", CYAN, True),
    ("02", "Human review", "Accept / reject / correct\nand add reviewer notes.", PALE_BLUE, False),
    ("03", "SFT dataset", "Group reviewed labels\nby chunk into JSONL.", PALE_BLUE, False),
    ("04", "QLoRA / SFT", "Train Qwen2.5-7B-Instruct,\nnot the AWQ model.", PALE_BLUE, False),
    ("05", "Before / after", "Compare precision, JSON stability,\nnew candidates and review effort.", PALE_GREEN, False),
]

start_x = 1.42
card_w = 3.12
gap = 0.39
for index, (step, title, body, fill, active) in enumerate(steps):
    x = start_x + index * (card_w + gap)
    add_card(roadmap_slide, x, 3.25, card_w, 4.65, step, title, body, fill, active)
    if index < len(steps) - 1:
        arrow = roadmap_slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(x + card_w + 0.05), Inches(5.15), Inches(0.3), Inches(0.65)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NAVY
        arrow.line.fill.background()

note = roadmap_slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(8.35), Inches(13.6), Inches(0.62)
)
note.fill.solid()
note.fill.fore_color.rgb = NAVY
note.line.fill.background()
add_text(
    roadmap_slide, 3.32, 8.4, 13.36, 0.48,
    "Training starts only after reviewed labels are sufficiently reliable (50-100 samples for a proof of concept).",
    17, WHITE, False
)

prs.save(OUTPUT)
print(OUTPUT)
